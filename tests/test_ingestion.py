from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from docx import Document as WordDocument
from pptx import Presentation

from document_analyst.config import AppSettings
from document_analyst.models import DocumentRecord, PageSpan
from document_analyst.services.ingestion import DocumentIngestor


class IngestionTests(unittest.TestCase):
    def test_discovery_is_case_insensitive_and_skips_symlinks(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            (root / "A.PDF").write_bytes(b"not parsed in discovery")
            (root / "notes.TXT").write_text("hello", encoding="utf-8")
            (root / "ignore.csv").write_text("no", encoding="utf-8")
            try:
                (root / "linked.txt").symlink_to(root / "notes.TXT")
            except OSError:
                pass
            paths = DocumentIngestor(AppSettings()).discover(str(root))
            self.assertEqual([item.name for item in paths], ["A.PDF", "notes.TXT"])

    def test_chunk_ids_are_fixed_size_and_deterministic(self) -> None:
        ingestor = DocumentIngestor(AppSettings(chunk_size=100, chunk_overlap=10))
        document = DocumentRecord(
            source_path="/very/long/" + "x" * 500, name="sample.txt", extension=".txt",
            size_bytes=20, text="First sentence. Second sentence.",
            page_spans=[PageSpan(1, "First sentence. Second sentence.")],
        )

        def embeddings(texts: list[str]) -> list[list[float]]:
            return [[1.0, 0.0] for _ in texts]

        chunks = ingestor.build_chunks([document], embeddings)
        self.assertTrue(chunks)
        self.assertLess(len(chunks[0].chunk_id), 80)
        self.assertEqual(chunks[0].chunk_id, ingestor.build_chunks([document], embeddings)[0].chunk_id)

    def test_chunking_progress_reports_each_document(self) -> None:
        ingestor = DocumentIngestor(AppSettings(chunk_size=100, chunk_overlap=10))
        documents = [
            DocumentRecord(
                source_path=f"/tmp/{name}", name=name, extension=".txt",
                size_bytes=10, text="A sentence.",
                page_spans=[PageSpan(1, "A sentence.")],
            )
            for name in ("first.txt", "second.txt")
        ]
        events: list[tuple[str, str, int, int, bool]] = []

        ingestor.build_chunks(
            documents,
            lambda texts: [[1.0, 0.0] for _ in texts],
            progress=lambda *event: events.append(event),
        )

        self.assertEqual(
            events,
            [
                ("chunking", "first.txt", 1, 2, False),
                ("chunking", "first.txt", 1, 2, True),
                ("chunking", "second.txt", 2, 2, False),
                ("chunking", "second.txt", 2, 2, True),
            ],
        )

    def test_docx_paragraphs_and_tables_are_extracted(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "report.docx"
            document = WordDocument()
            document.add_paragraph("Quarterly summary.")
            table = document.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "Revenue"
            table.cell(0, 1).text = "$42"
            document.save(path)

            loaded = DocumentIngestor(AppSettings()).load_document(path)

            self.assertIn("Quarterly summary.", loaded.text)
            self.assertIn("Revenue | $42", loaded.text)
            self.assertEqual(loaded.page_spans[0].page_number, 1)

    def test_pptx_text_and_slide_numbers_are_extracted(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "deck.pptx"
            presentation = Presentation()
            first = presentation.slides.add_slide(presentation.slide_layouts[1])
            first.shapes.title.text = "Overview"
            first.placeholders[1].text = "First slide body"
            second = presentation.slides.add_slide(presentation.slide_layouts[1])
            second.shapes.title.text = "Details"
            presentation.save(path)

            loaded = DocumentIngestor(AppSettings()).load_document(path)

            self.assertIn("Overview", loaded.text)
            self.assertIn("Details", loaded.text)
            self.assertEqual([span.page_number for span in loaded.page_spans], [1, 2])


if __name__ == "__main__":
    unittest.main()
