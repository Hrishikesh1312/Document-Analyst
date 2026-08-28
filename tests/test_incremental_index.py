from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from document_analyst.config import AppSettings
from document_analyst.services.index_manifest import IndexManifest
from document_analyst.services.rag import RagService


class IncrementalIndexTests(unittest.TestCase):
    def _service(self, directory: str) -> RagService:
        root = Path(directory)
        service = RagService(AppSettings(chroma_dir=str(root / "chroma")))
        service.manifest = IndexManifest(root / "manifest.json")
        service._embed_texts = lambda texts: [
            [1.0, float(index + 1)] for index, _text in enumerate(texts)
        ]
        return service

    def test_unchanged_changed_duplicate_and_removed_files(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            documents = Path(directory) / "documents"
            documents.mkdir()
            first = documents / "first.txt"
            first.write_text("Initial indexed content.", encoding="utf-8")
            service = self._service(directory)

            initial = service.index_documents(str(documents))
            unchanged = service.index_documents(str(documents))
            first.write_text("Changed indexed content.", encoding="utf-8")
            changed = service.index_documents(str(documents))
            duplicate = documents / "duplicate.txt"
            duplicate.write_text("Changed indexed content.", encoding="utf-8")
            duplicated = service.index_documents(str(documents))
            first.unlink()
            removed = service.index_documents(str(documents))

            self.assertEqual(initial.document_count, 1)
            self.assertEqual(unchanged.unchanged_count, 1)
            self.assertEqual(changed.document_count, 1)
            self.assertEqual(duplicated.duplicate_count, 1)
            self.assertEqual(removed.removed_count, 1)
            statuses = {item["document_name"]: item for item in service.document_statuses()}
            self.assertEqual(statuses["first.txt"]["status"], "removed")
            self.assertTrue(statuses["duplicate.txt"]["indexed_at"])

    def test_failed_file_can_be_retried(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            documents = Path(directory) / "documents"
            documents.mkdir()
            deck = documents / "deck.pptx"
            deck.write_bytes(b"not a presentation")
            service = self._service(directory)

            failed = service.index_documents(str(documents))
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Recovered presentation"
            presentation.save(deck)
            retried = service.index_documents(str(documents), retry_failed_only=True)

            self.assertEqual(failed.failed_count, 1)
            self.assertEqual(retried.document_count, 1)
            self.assertEqual(service.document_statuses()[0]["status"], "indexed")

    def test_cancellation_stops_before_indexing(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            documents = Path(directory) / "documents"
            documents.mkdir()
            (documents / "file.txt").write_text("content", encoding="utf-8")
            service = self._service(directory)

            result = service.index_documents(str(documents), should_cancel=lambda: True)

            self.assertTrue(result.cancelled)
            self.assertEqual(result.document_count, 0)
            self.assertEqual(service.stats()["chunks"], 0)


if __name__ == "__main__":
    unittest.main()
