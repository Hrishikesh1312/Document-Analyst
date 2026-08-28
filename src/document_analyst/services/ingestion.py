from __future__ import annotations

import re
import os
import hashlib
from pathlib import Path
from collections.abc import Callable

import pymupdf
from docx import Document as WordDocument
from pptx import Presentation

try:
    import pytesseract
except ImportError:  # pragma: no cover
    pytesseract = None

from document_analyst.config import AppSettings
from document_analyst.models import ChunkRecord, DocumentRecord, PageSpan

IndexProgress = Callable[[str, str, int, int, bool], None]


class DocumentIngestor:
    def __init__(self, settings: AppSettings) -> None:
        self.settings = settings
        self.supported_extensions = set(settings.supported_extensions)
        self.max_size_bytes = settings.max_file_size_mb * 1024 * 1024

    def discover(self, directory: str) -> list[Path]:
        if not directory.strip():
            raise ValueError("Choose a documents folder before building the index.")
        root = Path(directory).expanduser().resolve()
        if not root.is_dir():
            raise FileNotFoundError(f"Directory not found: {root}")
        discovered: list[Path] = []
        for current, directories, filenames in os.walk(root, followlinks=False):
            directories[:] = sorted(
                (name for name in directories if not (Path(current) / name).is_symlink()),
                key=str.casefold,
            )
            for name in sorted(filenames, key=str.casefold):
                path = Path(current) / name
                if not path.is_symlink() and path.suffix.lower() in self.supported_extensions:
                    discovered.append(path)
        return discovered

    def discover_legacy_presentations(self, directory: str) -> list[Path]:
        """Find unsupported binary PowerPoint files without following symlinks."""
        if not directory.strip():
            raise ValueError("Choose a documents folder before building the index.")
        root = Path(directory).expanduser().resolve()
        if not root.is_dir():
            raise FileNotFoundError(f"Directory not found: {root}")
        discovered: list[Path] = []
        for current, directories, filenames in os.walk(root, followlinks=False):
            directories[:] = sorted(
                (name for name in directories if not (Path(current) / name).is_symlink()),
                key=str.casefold,
            )
            for name in sorted(filenames, key=str.casefold):
                path = Path(current) / name
                if not path.is_symlink() and path.suffix.lower() == ".ppt":
                    discovered.append(path)
        return discovered

    def load_documents(
        self, directory: str, progress: IndexProgress | None = None
    ) -> tuple[list[DocumentRecord], list[str]]:
        documents: list[DocumentRecord] = []
        warnings: list[str] = []
        warned_about_ocr_runtime = False
        paths = self.discover(directory)
        total = len(paths)
        for index, path in enumerate(paths, start=1):
            if progress:
                progress("reading", path.name, index, total, False)
            try:
                if path.stat().st_size > self.max_size_bytes:
                    warnings.append(f"Skipped {path.name}: larger than {self.settings.max_file_size_mb}MB")
                else:
                    document = self.load_document(path)
                    if document.text:
                        documents.append(document)
                    else:
                        warnings.append(f"Skipped {path.name}: no readable text found")
            except Exception as exc:  # pragma: no cover - defensive runtime handling
                warnings.append(f"Skipped {path.name}: {exc}")
            if self.settings.enable_ocr and pytesseract is None and not warned_about_ocr_runtime:
                warnings.append("OCR is enabled, but the `pytesseract` Python package is not available.")
                warned_about_ocr_runtime = True
            if progress:
                progress("reading", path.name, index, total, True)
        return documents, warnings

    def load_document(self, path: Path) -> DocumentRecord:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            page_spans = self._read_pdf(path)
            text = "\n\n".join(page.text for page in page_spans)
        elif suffix == ".docx":
            page_spans = self._read_docx(path)
            text = "\n\n".join(span.text for span in page_spans)
        elif suffix == ".pptx":
            page_spans = self._read_pptx(path)
            text = "\n\n".join(span.text for span in page_spans)
        else:
            text = self._read_text(path)
            page_spans = [PageSpan(page_number=1, text=text)]

        return DocumentRecord(
            source_path=str(path.resolve()),
            name=path.name,
            extension=suffix,
            size_bytes=path.stat().st_size,
            text=text.strip(),
            page_spans=page_spans,
        )

    def build_chunks(
        self,
        documents: list[DocumentRecord],
        embed_texts: Callable[[list[str]], list[list[float]]],
        progress: IndexProgress | None = None,
    ) -> list[ChunkRecord]:
        chunks: list[ChunkRecord] = []
        total = len(documents)
        for index, document in enumerate(documents, start=1):
            if progress:
                progress("chunking", document.name, index, total, False)
            chunks.extend(self._chunk_document(document, embed_texts))
            if progress:
                progress("chunking", document.name, index, total, True)
        return chunks

    def _read_text(self, path: Path) -> str:
        for encoding in ("utf-8", "utf-8-sig", "latin-1"):
            try:
                return path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        raise UnicodeDecodeError("text", b"", 0, 1, "Unsupported text encoding")

    def _read_pdf(self, path: Path) -> list[PageSpan]:
        spans: list[PageSpan] = []
        with pymupdf.open(path) as doc:
            for index, page in enumerate(doc, start=1):
                text = page.get_text("text").strip()
                if self._should_run_ocr(text):
                    ocr_text = self._ocr_page(page).strip()
                    if len(ocr_text) > len(text):
                        text = ocr_text
                if text:
                    spans.append(PageSpan(page_number=index, text=text))
        return spans

    def _read_docx(self, path: Path) -> list[PageSpan]:
        document = WordDocument(path)
        blocks: list[str] = []
        blocks.extend(
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        )
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    blocks.append(" | ".join(cells))
        text = "\n".join(blocks).strip()
        return [PageSpan(page_number=1, text=text)] if text else []

    def _read_pptx(self, path: Path) -> list[PageSpan]:
        presentation = Presentation(path)
        spans: list[PageSpan] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            blocks: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if text:
                    blocks.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            blocks.append(" | ".join(cells))
            if blocks:
                spans.append(PageSpan(page_number=slide_number, text="\n".join(blocks)))
        return spans

    def _should_run_ocr(self, text: str) -> bool:
        return self.settings.enable_ocr and len(text.strip()) < self.settings.ocr_min_text_chars

    def _ocr_page(self, page: pymupdf.Page) -> str:
        if pytesseract is None:
            return ""
        from PIL import Image

        if self.settings.tesseract_cmd.strip():
            pytesseract.pytesseract.tesseract_cmd = self.settings.tesseract_cmd.strip()

        matrix = pymupdf.Matrix(self.settings.ocr_zoom, self.settings.ocr_zoom)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        mode = "RGB" if pix.n < 4 else "RGBA"
        image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
        if mode == "RGBA":
            image = image.convert("RGB")
        return pytesseract.image_to_string(image)

    def _chunk_document(
        self,
        document: DocumentRecord,
        embed_texts: Callable[[list[str]], list[list[float]]],
    ) -> list[ChunkRecord]:
        sentences, pages = self._split_sentences(document.page_spans)
        if not sentences:
            return []

        sentence_embeddings = embed_texts(sentences)
        chunks: list[ChunkRecord] = []
        current_sentences: list[str] = []
        current_pages: list[int] = []
        current_chars = 0
        sequence = 0

        for index, sentence in enumerate(sentences):
            similarity = 1.0
            if index > 0:
                previous = sentence_embeddings[index - 1]
                current = sentence_embeddings[index]
                similarity = sum(a * b for a, b in zip(previous, current))

            proposed_chars = current_chars + len(sentence) + 1
            semantic_break = similarity < self.settings.semantic_threshold and current_chars >= self.settings.chunk_overlap
            size_break = proposed_chars > self.settings.chunk_size and current_sentences

            if semantic_break or size_break:
                chunk_text = " ".join(current_sentences).strip()
                if chunk_text:
                    approx_page = current_pages[0] if current_pages else 1
                    chunks.append(
                        ChunkRecord(
                            chunk_id=self._chunk_id(document.source_path, sequence),
                            source_path=document.source_path,
                            document_name=document.name,
                            text=chunk_text,
                            sequence=sequence,
                            char_count=len(chunk_text),
                            approx_page=approx_page,
                            file_hash=document.file_hash,
                        )
                    )
                    sequence += 1

                overlap_text = self._tail_overlap(current_sentences)
                overlap_pages = current_pages[-len(overlap_text):] if overlap_text else []
                current_sentences = overlap_text
                current_pages = overlap_pages
                current_chars = sum(len(item) + 1 for item in current_sentences)

            current_sentences.append(sentence)
            current_pages.append(pages[index])
            current_chars += len(sentence) + 1

        final_text = " ".join(current_sentences).strip()
        if final_text:
            chunks.append(
                ChunkRecord(
                    chunk_id=self._chunk_id(document.source_path, sequence),
                    source_path=document.source_path,
                    document_name=document.name,
                    text=final_text,
                    sequence=sequence,
                    char_count=len(final_text),
                    approx_page=current_pages[0] if current_pages else 1,
                    file_hash=document.file_hash,
                )
            )
        return chunks

    def _split_sentences(self, page_spans: list[PageSpan]) -> tuple[list[str], list[int]]:
        sentences: list[str] = []
        pages: list[int] = []
        for span in page_spans:
            normalized = re.sub(r"\s+", " ", span.text).strip()
            parts = re.split(r"(?<=[.!?])\s+(?=[A-Z0-9\"'])", normalized)
            cleaned = [part.strip() for part in parts if part.strip()]
            sentences.extend(cleaned)
            pages.extend([span.page_number] * len(cleaned))
        return sentences, pages

    def _tail_overlap(self, sentences: list[str]) -> list[str]:
        if not sentences:
            return []
        overlap: list[str] = []
        total = 0
        for sentence in reversed(sentences):
            overlap.insert(0, sentence)
            total += len(sentence) + 1
            if total >= self.settings.chunk_overlap:
                break
        return overlap

    @staticmethod
    def _chunk_id(source_path: str, sequence: int) -> str:
        # Fixed-size IDs avoid backend limits and path separator differences.
        digest = hashlib.sha256(source_path.encode("utf-8", errors="surrogatepass")).hexdigest()
        return f"{digest}:{sequence}"
